import os
from reportlab.platypus import SimpleDocTemplate, Image
from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PILImage
import cv2
import imagehash
from PIL import Image
import os
import sys
import subprocess
import time

SIMULARITY_TOLERANCE = 8
BATCH_SIZE = 128

def print_percentage_slider(percentage):
    """Prints a percentage slider to the console.

    Args:
        percentage: The percentage to display.
    """

    sys.stdout.write("\r[")
    for i in range(100):
        if i <= percentage:
            sys.stdout.write("#")
        else:
            sys.stdout.write(" ")
    sys.stdout.write(f"] {percentage}%")
    sys.stdout.flush()

class FrameSelector():
    simularity_tolerance = 8
    batch_size = 128
    phash_dict = {}

    def __init__(self,video_path,output_dir) -> None:
        self.output_dir = output_dir
        self.video_path = video_path
        self.frame = None
        self.is_duplicate = False

    def get_frame_phash(self,frame):
        self.frame = frame
        frame_pil = Image.fromarray(cv2.cvtColor(self.frame, cv2.COLOR_BGR2GRAY))
        return imagehash.phash(frame_pil)
    
    def save_unique_frame(self,frame_count,phash):
        self.is_duplicate = False
        for existing_phash, existing_path in FrameSelector.phash_dict.items():
            if phash - existing_phash <= SIMULARITY_TOLERANCE:
                self.is_duplicate = True
                break
                        
        if not self.is_duplicate:
            # Add the pHash and image path to the dictionary
            frame_filename = os.path.join(self.output_dir, f"frame_{frame_count:04d}.jpg")
            FrameSelector.phash_dict[phash] = f"{self.video_path}/{frame_filename}"
            cv2.imwrite(frame_filename, self.frame)
            return frame_count+1
        return frame_count


class VideoProc():

    def __init__(self,video_path,output_dir) -> None:
        self.video_path = video_path
        self.output_dir = output_dir
        self.phash_dict = {}
        self.frameSelector = FrameSelector(video_path=video_path,output_dir=self.output_dir)
    
    def initializeProc(self):
        ImgToDoc.clean(self.output_dir)
        os.makedirs(self.output_dir, exist_ok=True)
        self.cap = cv2.VideoCapture(self.video_path)

    def process(self):
        print("Extracting Frames ...\n")

        if not self.cap.isOpened():
            print("\nError: Couldn't open the video file.")
            exit()
            
        frame_count = 0
        frame_index = 0
        percentage = 0
        prev_percentage = 0

        while True:
            ret, frame = self.cap.read()

            if not ret:
                break

            frame_rate = self.cap.get(cv2.CAP_PROP_FPS) 
            total_frames = int(self.cap.get(cv2.CAP_PROP_FRAME_COUNT))

            percentage = int((frame_index * 100)/total_frames)
            frame_index += 1

            time_seconds = frame_count / frame_rate

            frame_number = int(time_seconds * frame_rate)

            if prev_percentage != percentage:
               print_percentage_slider(percentage)
               prev_percentage = percentage

            
            frame_count = self.frameSelector.save_unique_frame(frame_count=frame_count,phash=self.frameSelector.get_frame_phash(frame))

    def close(self):
        self.cap.release()
        cv2.destroyAllWindows()
        print(f"\n Creating The presentation ...")




class ImgToDoc():
    def __init__(self,img_folder) -> None:
        self.image_folder = img_folder
        self.document_path = f"{self.image_folder}/presentation.pptx"

    def __enter__(self):
        pp = ImgToDoc.create_pptx_from_images(self.image_folder, self.document_path)
        print(f"\nPPTX created at {self.document_path}")
        return pp

    def __exit__(self,exc_type,exc_val,traceback):
        ImgToDoc.clean(self.image_folder)
        self.launch()


    @staticmethod
    def create_pptx_from_images(image_folder, pptx_path):
        prs = Presentation()

        for image_file in os.listdir(image_folder):
            if image_file.endswith(('.jpg', '.jpeg', '.png', '.gif')):  # Add more extensions if needed
                image_path = os.path.join(image_folder, image_file)
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout

                img_width, img_height = PILImage.open(image_path).size
                aspect_ratio = img_width / img_height

                # Determine the slide dimensions based on aspect ratio
                if aspect_ratio > 1:
                    slide_width = Inches(10)  # Landscape orientation
                    slide_height = slide_width / aspect_ratio
                else:
                    slide_height = Inches(7.5)  # Portrait orientation
                    slide_width = slide_height * aspect_ratio

                # Add the image to the slide, centered and larger
                left = (prs.slide_width - slide_width) / 2
                top = (prs.slide_height - slide_height) / 2
                pic = slide.shapes.add_picture(image_path, left, top, slide_width, slide_height)

        prs.save(pptx_path)
        return pptx_path
    
    @staticmethod
    def clean(image_folder):
        for filename in os.listdir(image_folder):
            if filename.endswith(('.jpg', '.jpeg', '.png', '.gif')):
                file_path = os.path.join(image_folder, filename)
                try:
                    os.remove(file_path)
                except Exception as e:
                    print(f"Error deleting {filename}: {e}")
    
    def launch(self):
        try:
            subprocess.Popen(['start', 'powerpnt', '/S', self.document_path], shell=True)
        except Exception as e:
            print(f"An error occurred: {e}")



def presExtractor():
    print(
        "Video from youtube url : 1\n"+
        "Video from local files : 2\n"
    )

    try:
       source = int(input())

       if(source == 1):
           print("\n not available atm ...")
           time.sleep(1)
           exit()
           
       else:
                    # Specify the path to the input video file
            video_path = input("Enter the path to the video : ")

            # Specify the output directory where the frames will be saved
            output_directory = input("Where are you willing to store the presentation : ")

            videop = VideoProc(video_path=video_path,output_dir=output_directory)
            videop.initializeProc()
            videop.process()

            with ImgToDoc(output_directory) as f:
                print("Finished ...")
           
    except Exception as e:
        pass


if __name__ == "__main__":
    presExtractor()